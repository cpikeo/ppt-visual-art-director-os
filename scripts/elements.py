# -*- coding: utf-8 -*-
"""
Layer 1 · Elements（元素层）

职责：把 text / shape / image 三类元素画成**原生可编辑**对象。
只依赖 primitives 与 RenderContext；不知道 spec 结构，不感知主题差异。
改本层不影响图表层与编排层。
"""
from __future__ import annotations

import tempfile
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

from primitives import (
    RenderContext, emu, pt, split_runs, is_cjk, estimate_lines,
    set_run_font, set_para_font, solid_fill, gradient_fill, stroke_color,
    ALIGN, ANCHOR, Emu,
)

SHAPE_TYPES = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "pie": MSO_SHAPE.PIE,
}


def normalize_fill(value):
    """将标准与历史 fill 表达统一为内部契约；不解析主题 token。"""
    if value is None:
        return {"type": "none"}
    if isinstance(value, str):
        return {"type": "solid", "color": value}
    if not isinstance(value, dict):
        raise ValueError(
            "Invalid fill format. Expected a color token or "
            "{type: solid|gradient|none, ...}."
        )
    if "type" in value:
        kind = str(value["type"]).lower()
        if kind == "none":
            return {"type": "none"}
        if kind == "solid":
            if not value.get("color"):
                raise ValueError("Invalid solid fill: missing required field 'color'.")
            return {"type": "solid", "color": value["color"], "opacity": value.get("opacity")}
        if kind == "gradient":
            stops = value.get("stops")
            if not isinstance(stops, list) or len(stops) < 2:
                raise ValueError("Invalid gradient fill: 'stops' must contain at least two stops.")
            return {"type": "gradient", "stops": stops,
                    "angle": float(value.get("angle", 90)),
                    "gradient_type": value.get("gradient_type", "linear")}
        raise ValueError(f"Invalid fill type {value['type']!r}; expected solid, gradient, or none.")
    # Legacy: {"color": "#fff", "opacity": 0.3}
    if "color" in value:
        return {"type": "solid", "color": value["color"], "opacity": value.get("opacity")}
    # Legacy: {"gradient": {"stops": [[position, color, opacity], ...]}}
    if "gradient" in value:
        g = value["gradient"]
        if not isinstance(g, dict):
            raise ValueError("Invalid legacy gradient: expected an object.")
        return {"type": "gradient", "stops": g.get("stops", []),
                "angle": float(g.get("angle", 90)), "gradient_type": "linear"}
    raise ValueError(
        "Invalid fill format. Received a mapping without 'type' or legacy 'color'. "
        "Expected {type: 'solid', color: '#fff', opacity: 0.3}."
    )


def _gradient_stops(stops, ctx):
    normalized = []
    for item in stops:
        if isinstance(item, dict):
            pos = item.get("position", item.get("pos"))
            col = item.get("color")
            alpha = item.get("opacity", item.get("alpha"))
        elif isinstance(item, (list, tuple)) and len(item) >= 2:
            pos, col = item[0], item[1]
            alpha = item[2] if len(item) > 2 else None
        else:
            raise ValueError("Invalid gradient stop; expected object or [position, color, opacity].")
        if pos is None or col is None:
            raise ValueError("Invalid gradient stop: position and color are required.")
        resolved = ctx.colors.get(col, col)
        if ctx.paint(resolved)[0] is None:
            raise ValueError(f"Invalid fill color token {col!r} in gradient stop.")
        normalized.append((float(pos), resolved, float(alpha) if alpha is not None else None))
    return normalized


def apply_fill(shape_or_fill, value, ctx: RenderContext, fallback=None):
    """Apply the canonical Fill Contract; legacy shapes are normalized first."""
    fill = getattr(shape_or_fill, "fill", shape_or_fill)
    spec = normalize_fill(value)
    if spec["type"] == "none":
        fill.background()
        return
    if spec["type"] == "gradient":
        gradient_fill(fill, _gradient_stops(spec["stops"], ctx), spec["angle"])
        return
    color, token_alpha = ctx.paint(spec["color"])
    if color is None:
        if fallback:
            color, token_alpha = ctx.paint(fallback)
        if color is None:
            raise ValueError(
                f"Invalid fill color {spec['color']!r}. Expected a theme role or #RGB/#RRGGBB token."
            )
    opacity = spec.get("opacity")
    solid_fill(fill, color, float(opacity) if opacity is not None else token_alpha)


# --------------------------------------------------------------------------
# text
# --------------------------------------------------------------------------
def _validate_text_contract(element: dict) -> None:
    """Fail early when a text element uses fields not consumed by this layer."""
    eid = element.get("id", "?")
    if "text" not in element:
        hint = "；检测到 content，请改用 text" if "content" in element else ""
        raise ValueError(
            f"TEXT_FIELD_MISSING: text 元素 '{eid}' 缺少必需字段 'text'{hint}"
        )
    if "content" in element:
        raise ValueError(
            f"TEXT_FIELD_INVALID: text 元素 '{eid}' 使用了未消费字段 'content'，请改用顶层字段 'text'"
        )
    if "style" in element:
        raise ValueError(
            f"TEXT_STYLE_INVALID: text 元素 '{eid}' 使用了未消费的嵌套字段 'style'；"
            "请将 size、color、bold、line_height 等属性放到元素顶层"
        )


def add_text(slide, element: dict, ctx: RenderContext) -> None:
    _validate_text_contract(element)
    x, y, w, h = ctx.bounds(element)
    cn, latin = ctx.families(element)
    color = ctx.text_color(element.get("color"))
    alpha = None
    if element.get("color") is not None:
        _, alpha = ctx.paint(element.get("color"))
    if alpha is None and element.get("opacity") is not None:
        alpha = float(element["opacity"])
    pad = float(element.get("padding", 0))
    wrap = bool(element.get("wrap", True))
    lh = float(element.get("line_height", 1.35))
    size = float(element.get("size", 18))
    size_pt = size * 0.75
    bold = bool(element.get("bold", False))
    italic = bool(element.get("italic", False))
    spacing = float(element.get("char_spacing", 0) or 0)
    uppercase = bool(element.get("uppercase", False))
    align = ALIGN.get(element.get("align"), PP_ALIGN.LEFT)

    tb = slide.shapes.add_textbox(Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h)))
    tb.name = str(element.get("id", "text"))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.vertical_anchor = ANCHOR.get(element.get("anchor"), ANCHOR[None])
    tf.margin_left = Emu(emu(pad))
    tf.margin_right = Emu(emu(pad))
    tf.margin_top = Emu(emu(pad))
    tf.margin_bottom = Emu(emu(pad))

    # 可选文本框底色（用于标签条 / 数据标牌）
    if element.get("fill"):
        apply_fill(tb, element["fill"], ctx)

    lines = str(element.get("text", "")).split("\n")
    total = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = pt(size * lh)
        if element.get("space_after") is not None:
            p.space_after = pt(float(element["space_after"]))
        if element.get("space_before") is not None:
            p.space_before = pt(float(element["space_before"]))
        if not line:
            total += 1
            continue
        total += estimate_lines(line, w - 2 * pad, size, wrap)
        for rv, iscjk in split_runs(line):
            run = p.add_run()
            run.text = rv
            set_run_font(run, cn if iscjk else latin, cn, size_pt, color,
                         bold, italic, spacing if not iscjk else None, alpha, uppercase)

    max_lines = element.get("max_lines")
    if isinstance(max_lines, int) and max_lines >= 1 and total > max_lines:
        ctx.warn(f"text '{element.get('id')}': 估算 {total} 行 > max_lines {max_lines}")
    need = total * size * lh
    if need > (h - 2 * pad) + 1:
        ctx.warn(f"text '{element.get('id')}': 估算高度 {need:.0f}px 超出文本框 "
                 f"{h - 2 * pad:.0f}px（建议增大 height 或减小 size）")


# --------------------------------------------------------------------------
# shape
# --------------------------------------------------------------------------
def shape_text(shape, element: dict, ctx: RenderContext) -> None:
    """形状内文字：默认垂直居中。"""
    tf = shape.text_frame
    tf.word_wrap = bool(element.get("text_wrap", True))
    tf.vertical_anchor = ANCHOR.get(element.get("text_anchor", "middle"), ANCHOR["middle"])
    cn, latin = ctx.families(element)
    color = ctx.text_color(element.get("text_color"))
    alpha = element.get("text_opacity")
    size = float(element.get("text_size", 16))
    spacing = float(element.get("char_spacing", 0) or 0)
    p = tf.paragraphs[0]
    p.text = str(element.get("text", ""))
    p.alignment = ALIGN.get(element.get("align"), PP_ALIGN.CENTER)
    p.line_spacing = pt(size * float(element.get("text_line_height", 1.25)))
    for run in p.runs:
        set_run_font(run, cn if is_cjk(run.text[:1] or " ") else latin, cn,
                     size * 0.75, color, bool(element.get("text_bold", False)),
                     False, spacing,
                     float(alpha) if alpha is not None else None,
                     bool(element.get("uppercase", False)))


def add_shape(slide, element: dict, ctx: RenderContext) -> None:
    x, y, w, h = ctx.bounds(element)
    kind = element.get("shape", "rect")
    stroke, stroke_alpha = ctx.paint(element.get("stroke"))
    if stroke_alpha is None and element.get("stroke_opacity") is not None:
        stroke_alpha = float(element["stroke_opacity"])
    sw = float(element.get("stroke_width", 1))

    if kind in ("line", "arrow"):
        c = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(emu(x)), Emu(emu(y)),
            Emu(emu(x + w)), Emu(emu(y + h)))
        c.name = str(element.get("id", "line"))
        if stroke is not None:
            stroke_color(c.line, stroke, stroke_alpha, sw)
        return

    if kind not in SHAPE_TYPES:
        ctx.warn(f"shape '{element.get('id')}': 未知 shape={kind!r}，回落到 rect")
        kind = "rect"
    s = slide.shapes.add_shape(
        SHAPE_TYPES[kind], Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h)))
    s.name = str(element.get("id", "shape"))

    fill_value = element.get("fill")
    if isinstance(fill_value, (int, float)) and not isinstance(fill_value, bool):
        # fill_opacity 与角色名搭配使用
        fill_value = element.get("fill_role", "primary")
        color, _ = ctx.paint(fill_value)
        solid_fill(s.fill, color, float(element["fill"]))
    elif fill_value is not None:
        apply_fill(s, fill_value, ctx)
        if element.get("fill_opacity") is not None and not isinstance(fill_value, dict):
            color, _ = ctx.paint(fill_value)
            if color is not None:
                solid_fill(s.fill, color, float(element["fill_opacity"]))
    else:
        s.fill.background()

    if stroke is not None:
        stroke_color(s.line, stroke, stroke_alpha, sw)
    else:
        s.line.fill.background()

    if element.get("text"):
        shape_text(s, element, ctx)


# --------------------------------------------------------------------------
# image（cover / contain 裁切）
# --------------------------------------------------------------------------
def _resolve_src(element: dict, base_path: str | None) -> Path:
    src = Path(element["src"])
    if src.is_absolute():
        return src
    root = Path(base_path).parent if base_path else Path.cwd()
    return (root / src).resolve()


def _fit_image(src: Path, w: float, h: float, fit: str, crop=None,
               bg: tuple | None = None) -> Path:
    from PIL import Image

    img = Image.open(src).convert("RGB")
    iw, ih = img.size
    if crop:
        l, t, r, b = crop
        img = img.crop((int(iw * l), int(ih * t), int(iw * (1 - r)), int(ih * (1 - b))))
        iw, ih = img.size
    sw, sh = (max(int(round(w)), 1), max(int(round(h)), 1))
    # 使用主题背景色填充 contain 留白，避免在暗色主题下出现白边。
    bg_fill = bg if bg is not None else (255, 255, 255)
    if fit == "contain":
        scale = min(sw / iw, sh / ih)
        nw, nh = max(int(round(iw * scale)), 1), max(int(round(ih * scale)), 1)
        img = img.resize((nw, nh), Image.LANCZOS)
        canvas = Image.new("RGB", (sw, sh), bg_fill)
        canvas.paste(img, ((sw - nw) // 2, (sh - nh) // 2))
        img = canvas
    else:  # cover
        scale = max(sw / iw, sh / ih)
        nw, nh = max(int(round(iw * scale)), 1), max(int(round(ih * scale)), 1)
        img = img.resize((nw, nh), Image.LANCZOS)
        img = img.crop(((nw - sw) // 2, (nh - sh) // 2,
                        (nw - sw) // 2 + sw, (nh - sh) // 2 + sh))
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    img.save(tmp.name, "PNG")
    return Path(tmp.name)


def add_image(slide, element: dict, ctx: RenderContext, base_path: str | None = None) -> None:
    x, y, w, h = ctx.bounds(element)
    src = _resolve_src(element, base_path)
    if not src.exists():
        ctx.warn(f"image '{element.get('id')}': 找不到文件 {src}")
        return

    fit = element.get("fit", "cover")
    crop = element.get("crop")
    temp_path = None
    path_to_use = src
    if fit in ("cover", "contain") or crop:
        try:
            # contain 留白填充主题背景色，保证暗色主题下无白边。
            bg_rgb = ctx.color("background")
            bg_tuple = tuple(bg_rgb) if bg_rgb is not None else None
            temp_path = _fit_image(src, w, h, fit, crop, bg=bg_tuple)
            path_to_use = temp_path
        except Exception as exc:
            ctx.warn(f"image '{element.get('id')}': 裁切失败，按原图嵌入（{exc}）")

    try:
        pic = slide.shapes.add_picture(str(path_to_use), Emu(emu(x)), Emu(emu(y)),
                                       Emu(emu(w)), Emu(emu(h)))
        pic.name = str(element.get("id", "image"))
    finally:
        if temp_path is not None:
            try:
                temp_path.unlink(missing_ok=True)
            except Exception:
                pass
