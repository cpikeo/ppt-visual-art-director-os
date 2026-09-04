# -*- coding: utf-8 -*-
"""
Layer 3 · Compiler（编排层）

职责：构造 RenderContext、遍历 spec、分发元素、收集诊断、落盘。
本层**不含任何绘制逻辑**，也不知道某个主题长什么样 —— 绘制在 elements / charts，
主题信息在 spec 里经 RenderContext 传递。

稳定对外契约（不随分层重构而改变）：
    compile_deck(spec, output_path) -> {"passed","slides","warnings","file_bytes"}

本文件是唯一推荐的调用入口。
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

_HERE = str(Path(__file__).resolve().parent)
if _HERE not in sys.path:
    sys.path.insert(0, _HERE)

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

from primitives import RenderContext, emu, DEFAULT_WIDTH, DEFAULT_HEIGHT  # noqa: E402
import elements  # noqa: E402
import charts  # noqa: E402

# 元素类型 → 绘制层函数（编排层只做映射，改绘制实现不影响此表）
DISPATCH = {
    "text": elements.add_text,
    "shape": elements.add_shape,
    "image": elements.add_image,
    "chart": charts.add_chart,
    "native_chart": charts.add_chart,
}


def compile_deck(spec: dict, output_path, checks: bool = True,
                 guard_rules: dict | None = None) -> dict:
    """
    把设计 spec 编译为原生可编辑 PPTX。

    spec（全部由调用方传入，引擎不持有任何主题）：
      canvas : {"width":1280,"height":720}            可选
      theme  : {"colors":{...},"fonts":{...},...}      设计身份
      slides : [{"id","background","elements":[...]}]  页面与元素

    checks=True（默认）时先做静态治理（engine/guard.py，OS 硬约束断言），
    静态问题以 [guard] 前缀并入 warnings；guard_rules 可配置治理阈值。
    返回契约不变：{"passed","slides","warnings","file_bytes"}；
    启用 checks 时追加 "guard"（治理明细），旧调用不受影响。
    """
    output_path = Path(output_path)
    canvas = dict(spec.get("canvas") or {"width": DEFAULT_WIDTH, "height": DEFAULT_HEIGHT})
    theme = spec.get("theme") or {}

    ctx = RenderContext(theme, canvas)

    guard = None
    if checks:
        from guard import check_spec
        guard = check_spec(spec, rules=guard_rules)
        for w in guard["warnings"]:
            ctx.warn(f"[guard] {w}")

    prs = Presentation()
    prs.slide_width = Emu(emu(canvas["width"]))
    prs.slide_height = Emu(emu(canvas["height"]))
    blank = prs.slide_layouts[6]

    slides = spec.get("slides") or []
    for si, slide_spec in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        bg = slide_spec.get("background", ctx.colors.get("background"))
        applied = False
        if bg is not None:
            # 支持纯色 / 带透明度 / 渐变背景（渐变用于营造方向性光影空间）
            try:
                elements.apply_fill(slide.background, bg, ctx)
                applied = True
            except Exception as exc:
                ctx.warn(f"slide[{si}] '{slide_spec.get('id')}': "
                         f"background={bg!r} 渲染失败（{exc}）")
        # 安全网：仅当主应用失败时才回落到纯色。
        # 注意必须用「成功标志」而不是检查 <p:cSld> 直接子元素——
        # 否则会把已正确写入 <p:bgPr> 的渐变误判为空并用纯色覆盖。
        if not applied:
            try:
                elements.apply_fill(slide.background, "background", ctx)
            except Exception:
                pass

        for ei, element in enumerate(slide_spec.get("elements") or []):
            fn = DISPATCH.get(str(element.get("type", "text")))
            if fn is None:
                ctx.warn(f"slide[{si}].elements[{ei}]: 未知 type="
                         f"{element.get('type')!r}，已跳过")
                continue
            try:
                if fn is elements.add_image:
                    fn(slide, element, ctx, str(output_path))
                else:
                    fn(slide, element, ctx)
            except Exception as exc:  # 只记录，不静默改稿
                ctx.warn(f"slide[{si}].elements[{ei}] ({element.get('id')}): {exc}")

    prs.save(str(output_path))
    report = {
        "passed": len(ctx.warnings) == 0,
        "slides": len(slides),
        "warnings": list(ctx.warnings),
        "file_bytes": output_path.stat().st_size,
    }
    if guard is not None:
        report["guard"] = {"score": guard["score"],
                           "checks": guard["checks"],
                           "passed": guard["passed"]}
    return report


# --------------------------------------------------------------------------
# 可选 CLI： python compiler.py <build_module.py> [output.pptx]
# build_module 需定义 build_spec() -> dict 或顶层常量 SPEC
# --------------------------------------------------------------------------
def main(argv):
    if len(argv) < 2:
        print("usage: python compiler.py <build_module.py> [output.pptx]")
        return 1
    mod_path = Path(argv[1])
    spec_mod = importlib.util.spec_from_file_location("buildmod", str(mod_path))
    mod = importlib.util.module_from_spec(spec_mod)
    spec_mod.loader.exec_module(mod)
    deck = mod.build_spec() if hasattr(mod, "build_spec") else getattr(mod, "SPEC", None)
    if deck is None:
        print("build module must define build_spec() or SPEC")
        return 1
    out = argv[2] if len(argv) > 2 else "deck.pptx"
    report = compile_deck(deck, out)
    print(report)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
