# -*- coding: utf-8 -*-
"""
Layer 2 · Charts（图表层）

职责：把 chart 元素渲染为**原生可编辑图表**或形状化图表。
只依赖 primitives 与 RenderContext；不知道 spec 结构，不感知主题差异。
改本层不影响元素层与编排层。
"""
from __future__ import annotations

import math

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN

from primitives import (
    RenderContext, emu, pt, set_para_font, solid_fill, ALIGN, ANCHOR, Emu, Pt,
)

NATIVE_CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "comparison_bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "trend": XL_CHART_TYPE.LINE,
    "single_trend_line": XL_CHART_TYPE.LINE,
    "area": XL_CHART_TYPE.AREA,
    "donut": XL_CHART_TYPE.DOUGHNUT,
    "donut_composition": XL_CHART_TYPE.DOUGHNUT,
    "pie": XL_CHART_TYPE.PIE,
}
SHAPE_CHARTS = {
    "process_flow", "timeline", "steps", "matrix",
    "waterfall", "architecture", "bubble",
    # 编辑级图表：轨道 + 细线 + 直接标注，避免通用图表外观
    "ranked_bar", "progress_bar", "stacked_bar", "big_number_row",
}


def _textbox(slide, name, x, y, w, h, text, size, color, ctx, element,
             align=PP_ALIGN.LEFT, bold=False, alpha=None):
    """图表内的可编辑文字块（垂直居中、零内边距）。"""
    cn, latin = ctx.families(element)
    tb = slide.shapes.add_textbox(Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h)))
    tb.name = name
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = ANCHOR["middle"]
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    set_para_font(p, latin, cn, size * 0.75, color, bold)
    return tb


def chart_colors(element: dict, ctx: RenderContext):
    primary = ctx.color(element.get("primary_color")
                        or ctx.theme.get("chart_primary", "accent"))
    if primary is None:
        primary = ctx.text_color("ink")
    secondary = ctx.color(element.get("secondary_color")
                          or ctx.theme.get("chart_secondary", "secondary"))
    if secondary is None:
        secondary = primary
    ink = ctx.text_color(element.get("ink_color"))
    muted = ctx.color(element.get("muted_color")
                      or ctx.theme.get("chart_muted", "secondary"))
    if muted is None:
        muted = ctx.text_color("secondary" if "secondary" in ctx.colors else None)
    return primary, secondary, ink, muted


def _rows(element: dict) -> list[dict]:
    """统一解析 data；保留无效值状态，禁止把坏数据静默当作 0。"""
    data = element.get("data", []) or []
    rows = []
    for i, it in enumerate(data):
        if isinstance(it, dict) and "label" in it:
            raw = it.get("value")
            missing = raw is None
            invalid = False
            try:
                num = float(raw) if raw is not None else 0.0
                invalid = not math.isfinite(num)
            except (TypeError, ValueError):
                num, invalid = 0.0, not missing
            if invalid:
                num = 0.0
            display = it.get("display")
            if display is None:
                display = raw if isinstance(raw, str) else (
                    int(num) if float(num).is_integer() else num)
            rows.append({"label": str(it["label"]), "value": num,
                         "display": display, "_index": i,
                         "missing": missing, "invalid": invalid})
        elif isinstance(it, (int, float)) and math.isfinite(float(it)):
            rows.append({"label": str(i), "value": float(it), "display": it,
                         "_index": i, "missing": False, "invalid": False})
    return rows


def _safe_index(value, default=0):
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _display(row, element, fallback=""):
    """Prefer caller-authored display strings; never invent a unit."""
    if row.get("missing"):
        return str(element.get("missing_label", "—"))
    value = row.get("display")
    if value is not None:
        return str(value)
    return fallback if fallback != "" else str(row.get("value", ""))


def _label(shape, text, element, ctx, color, default_size=14):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = ANCHOR["middle"]
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = PP_ALIGN.CENTER
    cn, latin = ctx.families(element)
    size = float(element.get("label_size", default_size))
    set_para_font(p, latin, cn, size * 0.75, color, False)


# --------------------------------------------------------------------------
# 原生图表
# --------------------------------------------------------------------------
def add_native_chart(slide, element: dict, ctx: RenderContext) -> None:
    x, y, w, h = ctx.bounds(element)
    kind = str(element.get("chart_kind") or element.get("kind", ""))
    cn, latin = ctx.families(element)
    primary, secondary, ink, muted = chart_colors(element, ctx)

    rows = _rows(element)
    if not rows:
        ctx.warn(f"chart '{element.get('id')}': 没有可渲染的数据行，已跳过")
        return
    invalid = [r for r in rows if r.get("invalid")]
    if invalid:
        ctx.warn(f"chart '{element.get('id')}': {len(invalid)} 个数据值无法解析为有限数字，已按 0 计算；请修正原始数据")
    data = CategoryChartData()
    data.categories = [r["label"] for r in rows]
    data.add_series(str(element.get("series_name", "")), [r["value"] for r in rows])
    gf = slide.shapes.add_chart(
        NATIVE_CHART_TYPES[kind], Emu(emu(x)), Emu(emu(y)),
        Emu(emu(w)), Emu(emu(h)), data)
    gf.name = str(element.get("id", "chart"))
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False

    try:
        # 低噪声：去网格线、去轴线，弱化刻度
        chart.value_axis.has_major_gridlines = False
        chart.value_axis.format.line.fill.background()
        chart.value_axis.tick_labels.font.size = Pt(8)
        chart.value_axis.tick_labels.font.color.rgb = muted
        chart.value_axis.tick_labels.font.name = latin
        chart.category_axis.format.line.fill.background()
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.color.rgb = ink
        chart.category_axis.tick_labels.font.name = latin
    except Exception:
        pass

    series = chart.series[0]
    try:
        if kind in ("line", "trend", "single_trend_line"):
            series.format.line.color.rgb = primary
            series.format.line.width = Pt(2.25)
            series.smooth = bool(element.get("smooth", False))
            # 折线图的唯一强调点：无标记则 highlight 形同虚设
            hl = _safe_index(element.get("highlight"), -1)
            if 0 <= hl < len(rows):
                try:
                    mk = series.points[hl].marker
                    mk.style = XL_MARKER_STYLE.CIRCLE
                    mk.size = 7
                    mk.format.fill.solid()
                    mk.format.fill.fore_color.rgb = secondary
                    mk.format.line.color.rgb = secondary
                except Exception:
                    pass
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = primary
            series.format.line.fill.background()
            try:
                chart.plots[0].gap_width = int(element.get("gap_width", 55))
            except Exception:
                pass

        hl = _safe_index(element.get("highlight"), -1)
        if 0 <= hl < len(rows):
            point = series.points[hl]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = secondary
            try:
                point.format.line.color.rgb = secondary
            except Exception:
                pass

        show_values = element.get("show_values", kind in ("comparison_bar", "bar", "column", "donut"))
        # Do not squeeze labels into a dense plot. The policy is explicit so the
        # caller can choose between removing redundant labels and failing QA.
        label_policy = str(element.get("label_collision_policy", "fail"))
        if show_values and len(rows) >= 6 and h < 190:
            if label_policy == "hide_redundant":
                show_values = False
                ctx.warn(f"chart '{element.get('id')}': 标签空间不足，按 hide_redundant 隐藏重复数值")
            else:
                ctx.warn(f"chart '{element.get('id')}': 标签空间不足；请拆图/减少类别，避免数值遮挡")
        if show_values:
            plot = chart.plots[0]
            plot.has_data_labels = True
            plot.data_labels.show_value = True
            plot.data_labels.font.size = Pt(9)
            plot.data_labels.font.color.rgb = ink
            plot.data_labels.font.name = latin
            nf = element.get("number_format")
            if nf:
                try:
                    plot.data_labels.number_format = str(nf)
                    plot.data_labels.number_format_is_linked = False
                except Exception:
                    pass
            try:
                plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass
            # 直接标注已承担读数职责时隐藏值轴刻度（OS §19.3：直接标注够用时
            # 不保留第二套读数通道），进一步降噪。
            try:
                chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
            except Exception:
                pass

        if kind in ("donut", "donut_composition", "pie"):
            if kind != "pie":
                chart.plots[0].hole_size = int(element.get("hole_size", 62))
            hl = _safe_index(element.get("highlight"), 0)
            for i, point in enumerate(series.points):
                point.format.fill.solid()
                # 构成图各扇区用系列色区分（原先非高亮点统一 secondary，
                # 多类构成会糊成同色块）；高亮点保持 accent 语义。
                point.format.fill.fore_color.rgb = (
                    ctx.color("accent") if i == int(hl) else ctx.series_color(i))
    except Exception:
        pass


# --------------------------------------------------------------------------
# 形状化图表（战略图 / 流程 / 时间轴等，保持原生可编辑）
# --------------------------------------------------------------------------
def add_shape_chart(slide, element: dict, ctx: RenderContext, kind: str) -> None:
    x, y, w, h = ctx.bounds(element)
    cn, latin = ctx.families(element)
    primary, secondary, ink, muted = chart_colors(element, ctx)
    rows = _rows(element)
    eid = str(element.get("id", kind))
    if not rows:
        ctx.warn(f"chart '{element.get('id')}': 没有可渲染的数据行，已跳过")
        return
    invalid = [r for r in rows if r.get("invalid")]
    if invalid:
        ctx.warn(f"chart '{element.get('id')}': {len(invalid)} 个数据值无法解析为有限数字，已按 0 计算；请修正原始数据")

    if kind == "process_flow":
        n = min(len(rows), 7)
        node_w = w / max(n, 1) * 0.72
        gap = w / max(n, 1) * 0.28
        for i, row in enumerate(rows[:7]):
            nx = x + i * (node_w + gap)
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(emu(nx)), Emu(emu(y + h * 0.30)),
                Emu(emu(node_w)), Emu(emu(h * 0.32)))
            s.name = f"{eid}__node_{i}"
            s.fill.solid()
            fill_rgb = primary if i == 0 else secondary
            s.fill.fore_color.rgb = fill_rgb
            s.line.fill.background()
            # 按实际填充色选文字色（比 on_primary/on_secondary 更稳——
            # chart_secondary 可能映射到 accent，深浅与 secondary 相反）
            text_color = ctx.auto_text_for(fill_rgb)
            _label(s, row["label"], element, ctx, text_color, 14)
            if i < n - 1:
                c = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, Emu(emu(nx + node_w)), Emu(emu(y + h * 0.46)),
                    Emu(emu(nx + node_w + gap)), Emu(emu(y + h * 0.46)))
                c.name = f"{eid}__edge_{i}"
                c.line.color.rgb = muted
                c.line.width = Emu(emu(1))
        return

    if kind == "timeline":
        n = min(len(rows), 7)
        axis_y = y + h * 0.52
        ax = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(emu(x)), Emu(emu(axis_y)),
            Emu(emu(x + w)), Emu(emu(axis_y)))
        ax.name = f"{eid}__axis"
        ax.line.color.rgb = muted
        ax.line.width = Emu(emu(1))
        for i, row in enumerate(rows[:7]):
            cx = x + w * (i + 0.5) / n
            d = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(emu(cx - 7)), Emu(emu(axis_y - 7)),
                Emu(emu(14)), Emu(emu(14)))
            d.name = f"{eid}__point_{i}"
            d.fill.solid()
            # OS §24：末节点为强调色（accent 是唯一强调语义，primary 通常已被标题占用）
            d.fill.fore_color.rgb = ctx.color("accent") if i == n - 1 else secondary
            d.line.fill.background()
            lb = slide.shapes.add_textbox(
                Emu(emu(cx - w / n / 2)), Emu(emu(axis_y + 16)),
                Emu(emu(w / n)), Emu(emu(h * 0.30)))
            lb.name = f"{eid}__label_{i}"
            p = lb.text_frame.paragraphs[0]
            p.text = row["label"]
            p.alignment = PP_ALIGN.CENTER
            set_para_font(p, latin, cn, float(element.get("label_size", 13)) * 0.75, ink, False)
        return

    if kind == "steps":
        n = min(len(rows), 6)
        col_w = w / max(n, 1)
        for i, row in enumerate(rows[:6]):
            cx = x + i * col_w
            nb = slide.shapes.add_textbox(
                Emu(emu(cx)), Emu(emu(y)), Emu(emu(col_w * 0.9)), Emu(emu(44)))
            nb.name = f"{eid}__n_{i}"
            p = nb.text_frame.paragraphs[0]
            p.text = str(i + 1)
            p.alignment = PP_ALIGN.LEFT
            set_para_font(p, latin, cn, float(element.get("num_size", 30)) * 0.75, secondary, True)
            tb = slide.shapes.add_textbox(
                Emu(emu(cx)), Emu(emu(y + 48)), Emu(emu(col_w * 0.9)), Emu(emu(h * 0.36)))
            tb.name = f"{eid}__t_{i}"
            p = tb.text_frame.paragraphs[0]
            p.text = row["label"]
            set_para_font(p, latin, cn, float(element.get("title_size", 18)) * 0.75, ink, True)
            db = slide.shapes.add_textbox(
                Emu(emu(cx)), Emu(emu(y + 48 + h * 0.36)),
                Emu(emu(col_w * 0.9)), Emu(emu(h * 0.5)))
            db.name = f"{eid}__d_{i}"
            p = db.text_frame.paragraphs[0]
            p.text = str(row.get("desc", ""))
            p.line_spacing = pt(float(element.get("desc_size", 14)) * 1.3)
            set_para_font(p, latin, cn, float(element.get("desc_size", 14)) * 0.75, muted, False)
            if i < n - 1:
                c = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, Emu(emu(cx + col_w * 0.9)), Emu(emu(y + 22)),
                    Emu(emu(cx + col_w)), Emu(emu(y + 22)))
                c.name = f"{eid}__c_{i}"
                c.line.color.rgb = muted
                c.line.width = Emu(emu(1))
        return

    if kind == "matrix":
        pts = element.get("points", []) or []
        pad_x, pad_y = w * 0.12, h * 0.12
        ox, oy = x + pad_x, y + h - pad_y
        pw, ph = w - pad_x * 1.25, h - pad_y * 1.35
        hx = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(emu(ox)), Emu(emu(oy)), Emu(emu(ox + pw)), Emu(emu(oy)))
        hx.name = f"{eid}__x"
        hx.line.color.rgb = muted
        hx.line.width = Emu(emu(1))
        vy = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(emu(ox)), Emu(emu(oy)), Emu(emu(ox)), Emu(emu(oy - ph)))
        vy.name = f"{eid}__y"
        vy.line.color.rgb = muted
        vy.line.width = Emu(emu(1))
        for i, p0 in enumerate(pts[:12]):
            nx = ox + float(p0["x"]) * pw
            ny = oy - float(p0["y"]) * ph
            d = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(emu(nx - 7)), Emu(emu(ny - 7)), Emu(emu(14)), Emu(emu(14)))
            d.name = f"{eid}__p_{i}"
            d.fill.solid()
            d.fill.fore_color.rgb = primary if p0.get("highlight") else secondary
            d.line.fill.background()
            lb = slide.shapes.add_textbox(
                Emu(emu(nx + 8)), Emu(emu(ny - 9)), Emu(emu(120)), Emu(emu(22)))
            lb.name = f"{eid}__l_{i}"
            p = lb.text_frame.paragraphs[0]
            p.text = str(p0.get("label", ""))
            set_para_font(p, latin, cn, 11 * 0.75, ink, False)
        return

    if kind == "waterfall":
        cum = 0.0
        starts, ends = [], []
        for r in rows:
            starts.append(cum)
            cum += r["value"]
            ends.append(cum)
        lo = min(0.0, *starts, *ends)
        hi = max(0.0, *starts, *ends)
        span = max(hi - lo, 1.0)
        plot_h = h * 0.62
        zero_y = y + h * 0.82 - (0.0 - lo) / span * plot_h
        ax = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(emu(x)), Emu(emu(zero_y)),
            Emu(emu(x + w)), Emu(emu(zero_y)))
        ax.name = f"{eid}__axis"
        ax.line.color.rgb = muted
        ax.line.width = Emu(emu(1))
        bw = w / max(len(rows), 1) * 0.46
        for i, r in enumerate(rows):
            left = x + w * (i + 0.5) / max(len(rows), 1) - bw / 2
            top = max(starts[i], ends[i])
            bot = min(starts[i], ends[i])
            ry = y + h * 0.82 - (top - lo) / span * plot_h
            rh = max(5, (top - bot) / span * plot_h)
            b = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(emu(left)), Emu(emu(ry)), Emu(emu(bw)), Emu(emu(rh)))
            b.name = f"{eid}__bar_{i}"
            b.fill.solid()
            b.fill.fore_color.rgb = primary if r["value"] >= 0 else secondary
            b.line.fill.background()
            lb = slide.shapes.add_textbox(
                Emu(emu(left - bw * 0.35)), Emu(emu(y + h * 0.82 + 10)),
                Emu(emu(bw * 1.7)), Emu(emu(24)))
            lb.name = f"{eid}__l_{i}"
            p = lb.text_frame.paragraphs[0]
            p.text = r["label"]
            p.alignment = PP_ALIGN.CENTER
            set_para_font(p, latin, cn, 11 * 0.75, muted, False)
        return

    if kind == "architecture":
        layers = element.get("layers", [])
        if not (1 <= len(layers) <= 3):
            layers = ["Layer 1", "Layer 2"]
        lh = h / len(layers) * 0.6
        for i, lt in enumerate(layers[:3]):
            ly = y + i * h / len(layers) + h / len(layers) * 0.2
            s = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(emu(x + w * 0.12)), Emu(emu(ly)),
                Emu(emu(w * 0.76)), Emu(emu(lh)))
            s.name = f"{eid}__l_{i}"
            fill_rgb = primary if i == 0 else secondary
            s.fill.solid()
            s.fill.fore_color.rgb = fill_rgb
            s.line.fill.background()
            _label(s, str(lt), element, ctx, ctx.auto_text_for(fill_rgb), 15)
        return

    if kind == "ranked_bar":
        # 编辑级排行条：细线轨道 + 圆头条 + 直接数值标注
        data = sorted(rows, key=lambda r: r["value"], reverse=True)
        n = min(len(data), 8)
        data = data[:n]
        label_w = w * float(element.get("label_ratio", 0.24))
        bar_x = x + label_w + 16
        value_w = float(element.get("value_width", 70))
        bar_right = x + w - value_w
        span = max(bar_right - bar_x, 10)
        maxv = max((r["value"] for r in data), default=1) or 1
        row_h = h / max(n, 1)
        bar_h = max(6, min(float(element.get("bar_height", 12)), row_h * 0.34))
        track_c, track_a = ctx.paint("track")
        hairline_c, hairline_a = ctx.paint("hairline")
        hl = _safe_index(element.get("highlight"), -1)
        for i, r in enumerate(data):
            cy = y + row_h * (i + 0.5)
            track = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(emu(bar_x)), Emu(emu(cy - 1)),
                Emu(emu(span)), Emu(emu(2)))
            track.name = f"{eid}__track_{i}"
            solid_fill(track.fill, track_c or hairline_c, track_a or hairline_a)
            track.line.fill.background()
            if r["value"] < 0:
                ctx.warn(f"chart '{element.get('id')}': ranked_bar 第 {i + 1} 行为负值，已不绘制负向长度")
            bw = span * max(r["value"], 0) / maxv
            color = (ctx.color("accent") if hl in (i, r.get("_index")) else primary)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(emu(bar_x)), Emu(emu(cy - bar_h / 2)),
                Emu(emu(max(bw, 4))), Emu(emu(bar_h)))
            bar.name = f"{eid}__bar_{i}"
            solid_fill(bar.fill, color)
            bar.line.fill.background()
            # 端点圆点：编辑级排行的标志（条末端锚点，强化"读到端点"）
            dot_d = bar_h + 6
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(emu(bar_x + max(bw, 4) - dot_d / 2)),
                Emu(emu(cy - dot_d / 2)), Emu(emu(dot_d)), Emu(emu(dot_d)))
            dot.name = f"{eid}__dot_{i}"
            solid_fill(dot.fill, color)
            dot.line.fill.background()
            _textbox(slide, f"{eid}__label_{i}", x, cy - row_h / 2, label_w, row_h,
                     r["label"], float(element.get("label_size", 13)), ink, ctx, element,
                     align=PP_ALIGN.RIGHT)
            _textbox(slide, f"{eid}__value_{i}", bar_x + bw + 8, cy - row_h / 2,
                     value_w, row_h, _display(r, element),
                     float(element.get("value_size", 12)), ink, ctx, element,
                     align=PP_ALIGN.LEFT, bold=True)
        return

    if kind == "progress_bar":
        # 进度条：轨道 + 填充 + 百分比直接标注
        n = min(len(rows), 6)
        data = rows[:n]
        row_h = h / max(n, 1)
        bar_h = max(6, min(float(element.get("bar_height", 10)), row_h * 0.22))
        label_w = w * float(element.get("label_ratio", 0.26))
        bar_x = x + label_w + 16
        value_w = float(element.get("value_width", 78))
        span = max(x + w - value_w - bar_x, 10)
        track_c, track_a = ctx.paint("track")
        for i, r in enumerate(data):
            cy = y + row_h * (i + 0.5)
            top = cy - bar_h / 2
            track = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Emu(emu(bar_x)), Emu(emu(top)),
                Emu(emu(span)), Emu(emu(bar_h)))
            track.name = f"{eid}__track_{i}"
            solid_fill(track.fill, track_c or ctx.color("faint"), track_a)
            track.line.fill.background()
            ceiling = float(r.get("max", element.get("max", 100))) or 100
            if r["value"] < 0:
                ctx.warn(f"chart '{element.get('id')}': progress_bar 第 {i + 1} 行为负值，已按 0% 绘制")
            ratio = max(0.0, min(1.0, r["value"] / ceiling))
            if ratio > 0:
                fill = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Emu(emu(bar_x)), Emu(emu(top)),
                    Emu(emu(max(span * ratio, 4))), Emu(emu(bar_h)))
                fill.name = f"{eid}__fill_{i}"
                solid_fill(fill.fill, ctx.series_color(i) if element.get(
                    "multi_color") else ctx.color("accent"))
                fill.line.fill.background()
            _textbox(slide, f"{eid}__label_{i}", x, cy - row_h / 2, label_w, row_h,
                     r["label"], float(element.get("label_size", 13)), ink, ctx, element,
                     align=PP_ALIGN.RIGHT)
            _textbox(slide, f"{eid}__value_{i}", x + w - value_w, cy - row_h / 2,
                     value_w, row_h, _display(r, element, f"{ratio * 100:.0f}%"),
                     float(element.get("value_size", 12)), ink, ctx, element,
                     align=PP_ALIGN.RIGHT, bold=True)
        return

    if kind == "stacked_bar":
        # 堆叠条：分段 + 直接图例
        negative_rows = [i for i, r in enumerate(rows) if r["value"] < 0]
        if negative_rows:
            ctx.warn(f"chart '{element.get('id')}': stacked_bar 含负值行 {negative_rows}，负值不绘制")
        total = sum(max(r["value"], 0) for r in rows) or 1
        bar_h = max(18, min(float(element.get("bar_height", 46)), h * 0.34))
        bar_y = y + (h - bar_h) / 2 - (h * 0.10 if element.get("legend") is not False else 0)
        cursor = x
        for i, r in enumerate(rows[:8]):
            seg = w * max(r["value"], 0) / total
            if seg <= 0:
                continue
            s = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Emu(emu(cursor)), Emu(emu(bar_y)),
                Emu(emu(seg)), Emu(emu(bar_h)))
            s.name = f"{eid}__seg_{i}"
            solid_fill(s.fill, ctx.series_color(i) if not element.get(
                "ramp") else ctx.ramp_color(i))
            s.line.fill.background()
            if element.get("show_values", True) and seg > 46:
                pct = r["value"] / total * 100
                _textbox(slide, f"{eid}__seg_label_{i}", cursor, bar_y, seg, bar_h,
                         _display(r, element, f"{pct:.0f}%"),
                         float(element.get("label_size", 12)),
                         ctx.auto_text_for(ctx.series_color(i) if not element.get("ramp") else ctx.ramp_color(i)),
                         ctx, element, align=PP_ALIGN.CENTER, bold=True)
            cursor += seg
        if element.get("legend") is not False:
            ly = bar_y + bar_h + 18
            lx = x
            for i, r in enumerate(rows[:8]):
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, Emu(emu(lx)), Emu(emu(ly + 4)), Emu(emu(10)), Emu(emu(10)))
                dot.name = f"{eid}__dot_{i}"
                solid_fill(dot.fill, ctx.series_color(i) if not element.get(
                    "ramp") else ctx.ramp_color(i))
                dot.line.fill.background()
                _textbox(slide, f"{eid}__legend_{i}", lx + 16, ly - 4, 150, 22,
                         r["label"], 11, ink, ctx, element)
                lx += 16 + 150 + 18
        return

    if kind == "big_number_row":
        # 指标行：大数字 + 说明 + 细线分隔
        n = min(len(rows) or len(element.get("items", [])), 5)
        items = rows[:5] or (element.get("items") or [])[:5]
        col_w = w / max(n, 1)
        for i, r in enumerate(items[:n]):
            cx = x + i * col_w
            vw = col_w - 22
            if i > 0:
                rule = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Emu(emu(cx - 11)), Emu(emu(y + h * 0.18)),
                    Emu(emu(1.5)), Emu(emu(h * 0.52)))
                rule.name = f"{eid}__rule_{i}"
                hc, ha = ctx.paint("hairline")
                solid_fill(rule.fill, hc or ctx.color("secondary"), ha or 0.28)
                rule.line.fill.background()
            _textbox(slide, f"{eid}__value_{i}", cx + 12, y + h * 0.14, vw, h * 0.40,
                     _display(r, element),
                     float(element.get("value_size", 52)),
                     ctx.color("primary"), ctx, element, bold=True)
            _textbox(slide, f"{eid}__label_{i}", cx + 12, y + h * 0.56, vw, h * 0.30,
                     r.get("label", ""), float(element.get("label_size", 14)),
                     ctx.color("muted"), ctx, element)
        return

    if kind == "bubble":
        negative_rows = [i for i, r in enumerate(rows) if r["value"] < 0]
        if negative_rows:
            ctx.warn(f"chart '{element.get('id')}': bubble 含负值行 {negative_rows}，负值按最小半径绘制")
        vals = [max(r["value"], 0.01) for r in rows]
        lo, hi = min(vals), max(vals)
        span = max(hi - lo, 0.01)
        cols = math.ceil(math.sqrt(max(len(rows), 1)))
        cw = w / cols
        chh = h / math.ceil(len(rows) / cols)
        for i, r in enumerate(rows[:12]):
            rad = 18 + 34 * (vals[i] - lo) / span
            cx = x + cw * (i % cols + 0.5)
            cy = y + chh * (i // cols + 0.5)
            d = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Emu(emu(cx - rad)), Emu(emu(cy - rad)),
                Emu(emu(rad * 2)), Emu(emu(rad * 2)))
            d.name = f"{eid}__b_{i}"
            fill_rgb = primary if i == _safe_index(element.get("highlight"), 0) else secondary
            d.fill.solid()
            d.fill.fore_color.rgb = fill_rgb
            d.line.fill.background()
            _label(d, r["label"], element, ctx, ctx.auto_text_for(fill_rgb), 11)
        return


# --------------------------------------------------------------------------
# 图表分发
# --------------------------------------------------------------------------
def add_chart(slide, element: dict, ctx: RenderContext) -> None:
    kind = str(element.get("chart_kind") or element.get("kind", ""))

    if kind in ("kpi", "executive_kpi", "big_number"):
        add_kpi(slide, element, ctx)
        return
    if kind in SHAPE_CHARTS:
        add_shape_chart(slide, element, ctx, kind)
        return
    if kind not in NATIVE_CHART_TYPES:
        ctx.warn(f"chart '{element.get('id')}': 不支持的 chart_kind {kind!r}，已跳过")
        return
    add_native_chart(slide, element, ctx)


def add_kpi(slide, element: dict, ctx: RenderContext) -> None:
    """大数字：拆成 value / label 两个可编辑文本框。"""
    x, y, w, h = ctx.bounds(element)
    cn, latin = ctx.families(element)
    primary, _secondary, _ink, muted = chart_colors(element, ctx)
    eid = str(element.get("id", "kpi"))
    align = element.get("align", "left")

    vb = slide.shapes.add_textbox(
        Emu(emu(x)), Emu(emu(y)), Emu(emu(w)), Emu(emu(h * 0.62)))
    vb.name = f"{eid}__value"
    p = vb.text_frame.paragraphs[0]
    p.text = str(element.get("value", ""))
    p.alignment = ALIGN.get(align, ALIGN[None])
    set_para_font(p, latin, cn, float(element.get("value_size", 56)) * 0.75, primary, True)

    lb = slide.shapes.add_textbox(
        Emu(emu(x)), Emu(emu(y + h * 0.66)), Emu(emu(w)), Emu(emu(h * 0.30)))
    lb.name = f"{eid}__label"
    p2 = lb.text_frame.paragraphs[0]
    p2.text = str(element.get("label", ""))
    p2.alignment = ALIGN.get(align, ALIGN[None])
    set_para_font(p2, latin, cn, float(element.get("label_size", 16)) * 0.75, muted, False)
